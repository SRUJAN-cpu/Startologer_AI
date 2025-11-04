import os
import sys
from pptx import Presentation
import PyPDF2
try:
    import docx  # python-docx
except ImportError:
    docx = None

# Add parent directory to path to import helpers
sys.path.append(os.path.dirname(os.path.dirname(__file__)))
from helpers.pdf_compressor import compress_pdf, should_compress_pdf
from services.online_pdf_compressor import compress_pdf_with_fallback

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_pdf(file_path):
    text = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text.append(page.extract_text())
    return "\n".join(text)

def extract_text_from_docx(file_path):
    if docx is None:
        raise ImportError("python-docx is not installed. Please install 'python-docx' to handle .docx files.")
    document = docx.Document(file_path)
    paragraphs = [p.text for p in document.paragraphs]
    return "\n".join(paragraphs)

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".pdf":
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)

        # For very large files (>30MB), use online compression
        if file_size_mb > 30:
            print(f"[TextExtractor] Large PDF ({file_size_mb:.1f}MB), using online compression...")
            compressed_path = compress_pdf_with_fallback(file_path)
            text = extract_text_from_pdf(compressed_path)
            # Clean up compressed file if it's different from original
            if compressed_path != file_path and os.path.exists(compressed_path):
                try:
                    os.remove(compressed_path)
                except Exception as e:
                    print(f"[TextExtractor] Failed to clean up compressed file: {e}")
            return text
        # For medium files (>10MB), use local compression
        elif should_compress_pdf(file_path, threshold_mb=10):
            print(f"[TextExtractor] PDF is large, compressing before extraction: {file_path}")
            compressed_path = compress_pdf(file_path)
            text = extract_text_from_pdf(compressed_path)
            # Clean up compressed file if it's different from original
            if compressed_path != file_path and os.path.exists(compressed_path):
                try:
                    os.remove(compressed_path)
                except Exception as e:
                    print(f"[TextExtractor] Failed to clean up compressed file: {e}")
            return text
        else:
            return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    else:
        raise ValueError("Unsupported file type: {}".format(ext))

